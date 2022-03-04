# -*- coding: utf-8 -*-

#
# Copyright (c) 2021 Sorcero, Inc.
#
# This file is part of Sorcero's Language Intelligence platform
# (see https://www.sorcero.com).
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU Lesser General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU Lesser General Public License for more details.
#
# You should have received a copy of the GNU Lesser General Public License
# along with this program. If not, see <http://www.gnu.org/licenses/>.
#

import os
import logging

from pydantic import BaseModel
from typing import Optional
from typing_extensions import Literal

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .. import documents
from .. import sources
from .base import BaseTransformer
from .tabular_document_create_md_passage import Transformer as MDTransformer
from .resource_create_text_document import Transformer as RTransformer
from ..utils import write_document_to_path

__logger__ = logging.getLogger("ingestum")
__script__ = os.path.basename(__file__).replace(".py", "")


class CropArea(BaseModel):
    left: float
    top: float
    right: float
    bottom: float


class Transformer(BaseTransformer):
    """
    Transforms a `PPTX` source into a `Text` document where the text document
    contains all human-readable text from the source.

    :param directory: Path to the directory where images will be extracted
    :type directory: str
    :param first_page: First page to be used
    :type first_page: int
    :param last_page: Last page to be used
    :type last_page: int
    :param crop: Dictionary with left, top, right and bottom coordinates to be
        included from the page, expressed in percentages. E.g. ``top=0.1`` and
        ``bottom=0.9``, means everything that comes before that first ten
        percent and that last ten percent will be excluded.
    :type crop: CropArea
    """

    class ArgumentsModel(BaseModel):
        first_page: Optional[int] = None
        last_page: Optional[int] = None
        crop: Optional[CropArea] = None
        directory: str

    class InputsModel(BaseModel):
        source: sources.PPTX

    class OutputsModel(BaseModel):
        document: documents.Text

    arguments: ArgumentsModel
    inputs: Optional[InputsModel]
    outputs: Optional[OutputsModel]

    type: Literal[__script__] = __script__

    def get_title(self, presentation):
        return presentation.core_properties.title

    @staticmethod
    def extract_resource_document(document, path, r_id):
        with open(path, "wb") as _file:
            _file.write(r_id.blob)

        return documents.Resource.new_from(None, content=path)

    @staticmethod
    def extract_tabular_document(table):
        _table = []

        row_count = len(table.rows)
        col_count = len(table.columns)
        for r in range(0, row_count):
            _row = []
            for c in range(0, col_count):
                cell = table.cell(r, c)
                paragraphs = cell.text_frame.paragraphs
                lines = []
                for paragraph in paragraphs:
                    _text = ""
                    for run in paragraph.runs:
                        _text = _text + run.text
                    lines.append(_text)
                _row.append("\n".join(lines))
            _table.append(_row)

        rows = len(_table)
        columns = len(_table[0]) if rows else 0

        return documents.Tabular.new_from(
            None, content=_table, rows=rows, columns=columns
        )

    def extract_image(self, document, r_id):
        name = "image_%06d" % self._images_count
        path = os.path.join(self.arguments.directory, name)
        document = self.extract_resource_document(document, path, r_id)

        return RTransformer.extract(document.content)

    def extract_table(self, table):
        document = self.extract_tabular_document(table)

        # dump it to the output directory
        name = "table_%06d.json" % self._tables_count
        path = os.path.join(self.arguments.directory, name)
        write_document_to_path(document, path)

        transformer = MDTransformer(format=None)
        return transformer.convert(document.content)

    def extract_table_from_shape(self, shape):
        if not self.filter(shape):
            return
        self._tables_count += 1
        table = shape.table
        self._lines += self.extract_table(table) + ".\n\n"

    def extract_image_from_shape(self, presentation, shape):
        if not self.filter(shape):
            return
        self._images_count += 1
        image = shape.image
        self._lines += self.extract_image(presentation, image) + ".\n\n"

    def extract_text_from_shape(self, shape):
        if self.filter(shape) and shape.text:
            self._lines += shape.text + ".\n\n"

    def iterate_group_shape(self, presentation, shape):
        for element in list(shape.shapes):
            if element.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.iterate_group_shape(presentation, element)
            else:
                if hasattr(element, "text"):
                    self.extract_text_from_shape(element)
                elif element.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self.extract_image_from_shape(presentation, element)
                elif element.has_table:
                    self.extract_table_from_shape(element)

    def filter(self, shape):
        # returns True when the shape will be extracted
        # returns False when the shape will be filtered because the lines go outside the crop area

        if self.arguments.crop is None:
            return True

        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape_left + shape.width
        shape_bottom = shape_top + shape.height

        # Lines that go outside the crop area
        if self.arguments.crop.left != -1 and (
            shape_left < self.arguments.crop.left * self._slide_width
        ):
            return False
        if self.arguments.crop.top != -1 and (
            shape_top < self.arguments.crop.top * self._slide_height
        ):
            return False
        if self.arguments.crop.right != -1 and (
            shape_right > self.arguments.crop.right * self._slide_width
        ):
            return False
        if self.arguments.crop.bottom != -1 and (
            shape_bottom > self.arguments.crop.bottom * self._slide_height
        ):
            return False

        return True

    def extract(self, presentation):
        self._tables_count = 0
        self._images_count = 0
        self._lines = ""
        self._slide_width = presentation.slide_width
        self._slide_height = presentation.slide_height

        first_page = (
            self.arguments.first_page - 1
            if self.arguments.first_page is not None and self.arguments.first_page > 0
            else 0
        )
        last_page = (
            self.arguments.last_page - 1
            if self.arguments.last_page is not None and self.arguments.last_page > 0
            else len(presentation.slides) - 1
        )

        self._slide_count = first_page

        for slide in list(presentation.slides)[first_page : last_page + 1]:
            self._slide_count += 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    self.extract_text_from_shape(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self.extract_image_from_shape(presentation, shape)
                elif shape.has_table:
                    self.extract_table_from_shape(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self.iterate_group_shape(presentation, shape)
            self._lines += f"__SLIDE__{self._slide_count}\n\n"

        content = self._lines.strip()

        return content

    def transform(self, source: sources.PPTX) -> documents.Text:
        super().transform(source=source)

        try:
            presentation = Presentation(source.path)
        except Exception as e:
            msg = "The file is not a PPTX."
            __logger__.error(
                msg,
                extra={
                    "props": {
                        "transformer": self.type,
                        "file": source.path,
                        "error": str(e),
                    }
                },
            )
            raise RuntimeError(msg)

        return documents.Text.new_from(
            source,
            content=self.extract(presentation),
            title=self.get_title(presentation),
        )
